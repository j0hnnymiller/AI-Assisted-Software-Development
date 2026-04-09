import unittest
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation

from scripts.generate_pptx import (
    add_centered_two_titles_slide,
    apply_markdown_formatting,
    build_presentation,
    mark_slide_hidden,
    parse_slide,
    split_marp_slides,
)


class SplitMarpSlidesTests(unittest.TestCase):
    def test_tab_indented_front_matter_is_stripped(self):
        md_content = """---
ai_generated: true
model: \"openai/gpt-5.3-codex@2026-03-28\"
prompt: |
\tcompare generate_pptx.py to regression manifest
\tand create new regression decks
marp: true
theme: default
paginate: true
---

# Architecture Truth || But Make It Memorable

::: notes
**Expected PPTX Rendering:**
- Layout: Centered Two Titles
:::

---

## Content After H1 Dividers

Body text.
"""

        slides = split_marp_slides(md_content)

        self.assertEqual(2, len(slides))
        self.assertTrue(slides[0].startswith("# Architecture Truth || But Make It Memorable"))
        self.assertNotIn("ai_generated:", slides[0])
        self.assertIn("**Expected PPTX Rendering:**", slides[0])

    def test_parse_slide_detects_hide_class_comment(self):
        slide_block = """<!-- _class: hide -->

## Placeholder for Instructor Bio

Body text
"""

        (
            title,
            body,
            _bg_image,
            _speaker_notes,
            _layout_name,
            _inline_images,
            _title_is_h1,
            _matter_of_fact_title,
            should_hide_slide,
        ) = parse_slide(slide_block)

        self.assertEqual("Placeholder for Instructor Bio", title)
        self.assertEqual("Body text", body)
        self.assertTrue(should_hide_slide)

    def test_mark_slide_hidden_sets_show_zero(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        mark_slide_hidden(slide)

        self.assertEqual("0", slide._element.get("show"))

    def test_regression_manifest_marks_hidden_slide_in_output(self):
        repo_root = Path(__file__).resolve().parent
        manifest_path = repo_root / "slides" / "manifests" / "regression-test-01.manifest.md"

        with TemporaryDirectory() as tmp_dir:
            output_path = Path(tmp_dir) / "regression-test-01-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            hidden_titles = [
                (slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else "")
                for slide in prs.slides
                if slide._element.get("show") == "0"
            ]

        self.assertIn("Hidden Slide Branch", hidden_titles)

    def test_regression_manifest_includes_table_layout_verification_slides(self):
        repo_root = Path(__file__).resolve().parent
        manifest_path = repo_root / "slides" / "manifests" / "regression-test-01.manifest.md"

        with TemporaryDirectory() as tmp_dir:
            output_path = Path(tmp_dir) / "regression-test-01-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            slide_titles = [
                slide.shapes.title.text.strip()
                for slide in prs.slides
                if slide.shapes.title and slide.shapes.title.text
            ]

        self.assertIn("Markdown Table Centering and Width", slide_titles)
        self.assertIn("Explicit Layout Table Branch", slide_titles)
        self.assertIn("Bold Formatting Branch", slide_titles)

    def test_h3_body_lines_render_as_styled_subheadings(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        text_frame = slide.placeholders[1].text_frame

        apply_markdown_formatting(text_frame, "### Important Subheading", paragraph=text_frame.paragraphs[0])

        paragraph = text_frame.paragraphs[0]
        self.assertEqual("Important Subheading", "".join(run.text for run in paragraph.runs))
        self.assertTrue(paragraph.runs)
        self.assertEqual(22, paragraph.runs[0].font.size.pt)
        self.assertTrue(paragraph.runs[0].font.bold)
        self.assertEqual(8, paragraph.space_before.pt)
        self.assertEqual(4, paragraph.space_after.pt)

    def test_centered_two_titles_subtitle_renders_without_bullet(self):
        prs = Presentation()

        add_centered_two_titles_slide(prs, "Matter of Fact", "Witty Subtitle")

        slide = prs.slides[0]
        subtitle_shape = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 1)
        paragraph = subtitle_shape.text_frame.paragraphs[0]

        self.assertEqual("Witty Subtitle", "".join(run.text for run in paragraph.runs))
        p_pr = paragraph._p.get_or_add_pPr()
        self.assertTrue(any(child.tag.endswith("}buNone") for child in p_pr))
        self.assertFalse(any(child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum") for child in p_pr))

    def test_standalone_double_asterisk_text_renders_bold(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        text_frame = slide.placeholders[1].text_frame

        apply_markdown_formatting(text_frame, "**Key Insight**", paragraph=text_frame.paragraphs[0])

        paragraph = text_frame.paragraphs[0]
        self.assertEqual("Key Insight", "".join(run.text for run in paragraph.runs))
        self.assertTrue(paragraph.runs)
        self.assertTrue(all(run.font.bold for run in paragraph.runs if run.text))

    def test_multiple_double_asterisk_segments_render_bold_without_literal_markers(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        text_frame = slide.placeholders[1].text_frame

        apply_markdown_formatting(
            text_frame,
            "This has **bold text** and **another bold span**.",
            paragraph=text_frame.paragraphs[0],
        )

        paragraph = text_frame.paragraphs[0]
        self.assertEqual(
            "This has bold text and another bold span.",
            "".join(run.text for run in paragraph.runs),
        )
        bold_runs = [run for run in paragraph.runs if run.font.bold]
        self.assertEqual(["bold text", "another bold span"], [run.text for run in bold_runs])

    def test_markdown_bold_in_title_placeholder_renders_as_bold_run(self):
        with TemporaryDirectory() as tmp_dir:
            repo_root = Path(tmp_dir)
            slides_dir = repo_root / "slides"
            marp_dir = slides_dir / "marp"
            marp_dir.mkdir(parents=True)

            slide_path = marp_dir / "bold-title.deck.md"
            slide_path.write_text(
                "---\nmarp: true\ntheme: default\npaginate: true\n---\n\n"
                "## **Bold Title Branch**\n\n"
                "Normal body text.\n",
                encoding="utf-8",
            )

            manifest_path = slides_dir / "bold-title.manifest.md"
            manifest_path.write_text(
                "sections:\n"
                "  - name: Intro\n"
                "    decks:\n"
                "      - slides/marp/bold-title.deck.md\n",
                encoding="utf-8",
            )

            output_path = slides_dir / "bold-title-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            title_paragraph = prs.slides[0].shapes.title.text_frame.paragraphs[0]

            self.assertEqual("Bold Title Branch", "".join(run.text for run in title_paragraph.runs))
            self.assertTrue(title_paragraph.runs)
            self.assertTrue(all(run.font.bold for run in title_paragraph.runs if run.text))

    def test_markdown_bold_in_table_cell_renders_as_bold_run(self):
        with TemporaryDirectory() as tmp_dir:
            repo_root = Path(tmp_dir)
            slides_dir = repo_root / "slides"
            marp_dir = slides_dir / "marp"
            marp_dir.mkdir(parents=True)

            slide_path = marp_dir / "bold-table.deck.md"
            slide_path.write_text(
                "---\nmarp: true\ntheme: default\npaginate: true\n---\n\n"
                "## Bold Table Cells\n\n"
                "| Case | Value |\n"
                "| ---- | ----- |\n"
                "| title | **Bold Cell** |\n",
                encoding="utf-8",
            )

            manifest_path = slides_dir / "bold-table.manifest.md"
            manifest_path.write_text(
                "sections:\n"
                "  - name: Intro\n"
                "    decks:\n"
                "      - slides/marp/bold-table.deck.md\n",
                encoding="utf-8",
            )

            output_path = slides_dir / "bold-table-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            table_shape = next(shape for shape in prs.slides[0].shapes if getattr(shape, "has_table", False))
            paragraph = table_shape.table.cell(1, 1).text_frame.paragraphs[0]

            self.assertEqual("Bold Cell", "".join(run.text for run in paragraph.runs))
            self.assertTrue(paragraph.runs)
            self.assertTrue(any(run.font.bold for run in paragraph.runs if run.text))

    def test_explicit_layout_table_still_renders_as_powerpoint_table(self):
        with TemporaryDirectory() as tmp_dir:
            repo_root = Path(tmp_dir)
            slides_dir = repo_root / "slides"
            marp_dir = slides_dir / "marp"
            marp_dir.mkdir(parents=True)

            slide_path = marp_dir / "table.deck.md"
            slide_path.write_text(
                "---\nmarp: true\ntheme: default\npaginate: true\n---\n\n"
                "<!-- layout: Two Content -->\n\n"
                "## Copilot Plan Overview\n\n"
                "| Feature | Individual | Business | Enterprise |\n"
                "| ------- | ---------- | -------- | ---------- |\n"
                "| Price | $10/mo | $19/user/mo | $39/user/mo |\n",
                encoding="utf-8",
            )

            manifest_path = slides_dir / "table.manifest.md"
            manifest_path.write_text(
                "sections:\n"
                "  - name: Intro\n"
                "    decks:\n"
                "      - slides/marp/table.deck.md\n",
                encoding="utf-8",
            )

            output_path = slides_dir / "table-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            self.assertTrue(any(getattr(shape, "has_table", False) for shape in prs.slides[0].shapes))

    def test_table_renders_centered_at_eighty_percent_slide_width(self):
        with TemporaryDirectory() as tmp_dir:
            repo_root = Path(tmp_dir)
            slides_dir = repo_root / "slides"
            marp_dir = slides_dir / "marp"
            marp_dir.mkdir(parents=True)

            slide_path = marp_dir / "table-layout.deck.md"
            slide_path.write_text(
                "---\nmarp: true\ntheme: default\npaginate: true\n---\n\n"
                "## Comparison Matrix\n\n"
                "| Aspect | Instruction Files | Prompt Files | Custom Agents |\n"
                "| ------ | ----------------- | ------------ | ------------- |\n"
                "| Scope | Repository-wide | Single task | Conversational |\n",
                encoding="utf-8",
            )

            manifest_path = slides_dir / "table-layout.manifest.md"
            manifest_path.write_text(
                "sections:\n"
                "  - name: Intro\n"
                "    decks:\n"
                "      - slides/marp/table-layout.deck.md\n",
                encoding="utf-8",
            )

            output_path = slides_dir / "table-layout-draft.pptx"
            build_presentation(manifest_path, output_path)

            prs = Presentation(output_path)
            table_shape = next(shape for shape in prs.slides[0].shapes if getattr(shape, "has_table", False))
            expected_width = int(prs.slide_width * 0.8)
            expected_left = int((prs.slide_width - expected_width) / 2)

            self.assertEqual(expected_width, table_shape.width)
            self.assertEqual(expected_left, table_shape.left)


if __name__ == "__main__":
    unittest.main()
