#!/usr/bin/env python3
"""
Convert PPTX to Markdown using python-pptx.
Extracts text (with shape context), speaker notes, and images from each slide.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_text_from_shape(shape):
    """Return (is_title, paragraphs_list) for a text-bearing shape."""
    if not shape.has_text_frame:
        return False, []

    is_title = shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "placeholder_format") and shape.placeholder_format is not None and shape.placeholder_format.idx in (0, 1)

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level  # 0 = top level, 1+ = sub-bullets
        paragraphs.append((level, text))

    return is_title, paragraphs


def extract_slide_data(slide, slide_num, images_dir, pptx_stem):
    """
    Extract all content from a slide.
    Returns a dict with keys: title, shapes, notes, images.
    """
    data = {
        "title": None,
        "shapes": [],   # list of (level, text) tuples
        "notes": None,
        "images": [],   # list of saved image relative paths
    }

    # --- Text shapes ---
    for shape in slide.shapes:
        if shape.has_text_frame:
            is_title, paragraphs = get_text_from_shape(shape)
            if is_title and paragraphs and data["title"] is None:
                data["title"] = paragraphs[0][1]
                # remaining paragraphs of title shape go to body
                for item in paragraphs[1:]:
                    data["shapes"].append(item)
            else:
                for item in paragraphs:
                    data["shapes"].append(item)

        # --- Images / picture shapes ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                ext = image.ext  # e.g. 'png', 'jpeg'
                img_filename = f"{pptx_stem}_slide{slide_num:02d}_{shape.shape_id}.{ext}"
                img_path = images_dir / img_filename
                img_path.write_bytes(image.blob)
                data["images"].append(img_filename)
                print(f"  Saved image: {img_filename}")
            except Exception as e:
                print(f"  Warning: could not save image on slide {slide_num}: {e}")

    # --- Speaker notes ---
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            data["notes"] = notes_text

    return data


def format_paragraphs(paragraphs):
    """Convert (level, text) list to indented Markdown bullet lines."""
    lines = []
    for level, text in paragraphs:
        indent = "  " * level
        # Detect if text already starts with a list marker
        if text.startswith(("•", "○", "▪", "-", "*", "–")):
            lines.append(f"{indent}- {text[1:].strip()}")
        else:
            lines.append(f"{indent}- {text}")
    return lines


def convert_pptx_to_md(pptx_path, md_path):
    """Convert a PPTX file to Markdown, saving images alongside the output."""
    pptx_path = Path(pptx_path)
    md_path = Path(md_path)

    print(f"Converting: {pptx_path}")

    if not pptx_path.exists():
        print(f"Error: Input file not found: {pptx_path}")
        return False

    # Create output directory and images sub-folder
    md_path.parent.mkdir(parents=True, exist_ok=True)
    images_dir = md_path.parent / "images"
    images_dir.mkdir(exist_ok=True)

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f"Error opening PPTX: {e}")
        return False

    total_slides = len(prs.slides)
    print(f"Found {total_slides} slides")

    pptx_stem = pptx_path.stem.replace(" ", "_")

    try:
        with open(md_path, "w", encoding="utf-8") as f:
            # Document title from file name
            doc_title = pptx_path.stem.lstrip("_").replace("-", " ").replace("_", " ")
            f.write(f"# {doc_title}\n\n")
            f.write(f"*Converted from: `{pptx_path.name}`*\n\n---\n\n")

            for slide_idx, slide in enumerate(prs.slides, start=1):
                data = extract_slide_data(slide, slide_idx, images_dir, pptx_stem)

                # Slide heading
                if data["title"]:
                    f.write(f"## Slide {slide_idx}: {data['title']}\n\n")
                else:
                    f.write(f"## Slide {slide_idx}\n\n")

                # Body text
                if data["shapes"]:
                    body_lines = format_paragraphs(data["shapes"])
                    f.write("\n".join(body_lines))
                    f.write("\n\n")

                # Embedded images
                for img_filename in data["images"]:
                    f.write(f"![Slide {slide_idx} image](images/{img_filename})\n\n")

                # Speaker notes — pandoc fenced div format
                if data["notes"]:
                    f.write("::: notes\n")
                    f.write(data["notes"])
                    f.write("\n:::\n\n")

                f.write("---\n\n")

                print(f"  Slide {slide_idx}: {len(data['shapes'])} text items, "
                      f"{len(data['images'])} image(s), "
                      f"{'notes' if data['notes'] else 'no notes'}")

        print(f"\nSuccessfully created: {md_path}")
        print(f"Converted {total_slides} slides")
        return True

    except Exception as e:
        print(f"Error writing Markdown: {e}")
        return False


if __name__ == "__main__":
    # Default paths
    pptx_file = "Slides/individual-slides/Fundementals-Agenda.pptx"
    md_file = "Slides/individual-slides/Fundementals-Agenda.md"

    # Allow command line arguments
    if len(sys.argv) >= 2:
        pptx_file = sys.argv[1]
    if len(sys.argv) >= 3:
        md_file = sys.argv[2]

    success = convert_pptx_to_md(pptx_file, md_file)
    sys.exit(0 if success else 1)
