#!/usr/bin/env python3
"""
Extract PPTX to Marp Markdown format with speaker notes.
Requires: pip install python-pptx
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

TITLE_PLACEHOLDER_TYPES = {1, 13, 15}   # TITLE, CENTER_TITLE, and legacy values
SUBTITLE_PLACEHOLDER_TYPES = {2, 12}    # SUBTITLE


def get_placeholder_type(shape):
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type
    except Exception:
        pass
    return None


def extract_text_from_shape(shape):
    """Extract text from a non-title shape, including tables."""
    text_parts = []
    try:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    level = paragraph.level
                    if level > 0:
                        indent = "  " * (level - 1)
                        text_parts.append(f"{indent}- {text}")
                    else:
                        text_parts.append(text)
        elif shape.shape_type == 19:  # TABLE
            try:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    text_parts.append(" | ".join(rows[0]))
                    text_parts.append(" | ".join(["---"] * len(rows[0])))
                    for row in rows[1:]:
                        text_parts.append(" | ".join(row))
            except Exception:
                pass
    except Exception:
        pass
    return text_parts


def convert_pptx_to_marp(pptx_path, output_path):
    """Convert PPTX to Marp Markdown with ::: notes speaker notes."""
    print(f"Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)

    lines = ["---", "marp: true", "theme: default", "paginate: true", "---", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {slide_num}...")

        title_text = None
        body_lines = []

        for shape in slide.shapes:
            ph_type = get_placeholder_type(shape)
            if ph_type in TITLE_PLACEHOLDER_TYPES or ph_type in SUBTITLE_PLACEHOLDER_TYPES:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        if title_text is None:
                            title_text = t
                        else:
                            body_lines.append(t)
            else:
                body_lines.extend(extract_text_from_shape(shape))

        if title_text:
            lines.append(f"## {title_text}")
        if body_lines:
            lines.append("")
            lines.extend(body_lines)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("")
                lines.append("::: notes")
                for note_line in notes_text.splitlines():
                    lines.append(note_line.strip())
                lines.append(":::")

        lines.append("")
        lines.append("---")
        lines.append("")

    print(f"Writing output to: {output_path}")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"Successfully converted {len(prs.slides)} slides")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_to_marp.py <input.pptx> [output.md]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) >= 3 else str(Path(input_path).with_suffix(".md"))

    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    convert_pptx_to_marp(input_path, output_path)
