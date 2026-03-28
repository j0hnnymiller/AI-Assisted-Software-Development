"""
generate_pptx.py -- Build a PPTX from a YAML manifest using python-pptx.
Requires: pip install python-pptx pyyaml

Usage:
    python generate_pptx.py <yaml_path> <output_pptx_path>

🔒 CRITICAL INVARIANTS (DO NOT BREAK):

1. FIRST SECTION EXCEPTION:
    The first section (idx=0) in the YAML manifest MUST NOT receive the injected
    module list slide. Only sections with idx > 0 receive that injected slide.

    Rationale: First section typically contains welcome/title slides. Adding
    navigation before it creates an awkward opening that frontloads course
    structure.

    Implementation: The section loop MUST use enumerate() and check: if idx > 0
    before adding add_module_list_slide().

2. SPEAKER NOTES ON ALL SLIDES:
   Every slide (injected and content) MUST have speaker notes:
   - Injected slides: Notes explaining auto-generation and purpose
   - Content slides: Notes showing source file path (e.g., "Source: slides\\...")

   Rationale: Provides context in PowerPoint for instructors and maintainers.

   Implementation: All add_*_slide() functions MUST accept 'note' parameter and call
   set_slide_notes(slide, note) when note is provided.

3. MARKDOWN LINK RENDERING:
   Markdown links [text](url) are processed according to these rules:
   - If link text == link target: render as just "text"
   - If link text != link target: render as "text: url"

   Rationale: Reduces redundancy for self-describing URLs while preserving
   clarity when link text differs from target.

   Implementation: process_markdown_links() function in parse_slide()

4. MARKDOWN BOLD FORMATTING:
   Markdown bold syntax **text** is parsed and rendered as actual bold formatting in PPTX.

   Rationale: Ensures emphasis and visual hierarchy from source markdown is preserved
   in the PowerPoint output without manual reformatting.

   Implementation: apply_markdown_formatting() function parses **bold** syntax and
   creates text runs with font.bold = True. Used by add_title_content_slide() when
   adding body content to slides.

   Example: "**Principal Engineer**" renders as bold text, not literal asterisks.

5. MARKDOWN BLOCKQUOTE RENDERING:
   Blockquote lines (starting with '> ') are rendered as italic text in PPTX.

   Rationale: PPTX has no native blockquote element. Italic is the closest visual
   equivalent and distinguishes quoted material from normal body text.

   Implementation: apply_markdown_formatting() detects lines starting with '> ',
   strips the prefix, and wraps the content in italic runs.

   Example: "> 'Programming hasn't changed'" renders as italic text without the '> '.

6. CENTERED TITLE SLIDES FOR H1 HEADINGS:
   Any slide block whose title comes from a bare `# H1` heading, has no body
   content, and carries no explicit layout directive is rendered as a
   Centered Title (Section Header layout) slide — EXCEPT for the very first
   H1 encountered in the very first deck of the manifest, which is kept as a
   normal Title Only slide (or whatever its explicit layout dictates).

   Rationale: `# H1` headings in deck files act as deck-cover dividers.
   Rendering them with the Section Header layout gives a visually distinct
   centred title that signals a new topic without adding body content.

   Implementation: parse_slide() returns title_is_h1=True when the title was
   sourced from `# `.  build_presentation() tracks global_deck_idx and
   first_h1_seen to apply add_centered_title_slide() for all H1-only blocks
   after the first one.

See .github/instructions/slide-pipeline.instructions.md for full specification.
Use the verification checklist in that file before committing changes.
"""
import argparse
import hashlib
import re
import shutil
import subprocess
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

LAYOUT_TITLE_SLIDE = 0  # Title Slide layout
LAYOUT_TITLE_CONTENT = 1  # adjust index if template differs
LAYOUT_SECTION_HEADER = 2  # adjust index if template differs
LAYOUT_TWO_COLUMN = 3  # adjust index if template differs
LAYOUT_TITLE_ONLY = 5  # adjust index if template differs
LAYOUT_CENTERED_TWO_TITLES = 11  # Centered Two Titles layout (title + subtitle, both centred)

ALLOWED_SLIDE_EXTENSIONS = {".md", ".markdown"}
MERMAID_FENCE_PATTERN = re.compile(r"```mermaid\s*\n(.*?)\n```", re.DOTALL | re.IGNORECASE)


def detect_mermaid_cli() -> list[str] | None:
    """Return the Mermaid CLI invocation command, or None if unavailable."""
    candidates: list[list[str]] = []
    if shutil.which("mmdc"):
        candidates.append(["mmdc"])
    if shutil.which("npx"):
        candidates.append(["npx", "-y", "@mermaid-js/mermaid-cli"])

    import os
    for cmd in candidates:
        try:
            subprocess.run(
                [*cmd, "--version"],
                check=True,
                capture_output=True,
                text=True,
                shell=(os.name == "nt")
            )
            return cmd
        except Exception:
            continue

    return None


def extract_mermaid_blocks(body: str) -> tuple[str, list[str]]:
    """Extract Mermaid fenced blocks and return (body_without_mermaid, blocks)."""
    matches = MERMAID_FENCE_PATTERN.findall(body)
    body_without_mermaid = MERMAID_FENCE_PATTERN.sub("", body)
    return body_without_mermaid.strip(), [m.strip() for m in matches if m.strip()]


def render_mermaid_png(
    mermaid_code: str,
    cache_dir: Path,
    mermaid_cli: list[str],
    width: int = 1600,
    height: int = 900,
) -> Path | None:
    """Render Mermaid code to PNG and return output path, or None on failure."""
    digest = hashlib.sha1(mermaid_code.encode("utf-8")).hexdigest()
    cache_dir.mkdir(parents=True, exist_ok=True)

    input_path = cache_dir / f"{digest}.mmd"
    output_path = cache_dir / f"{digest}.png"

    if output_path.exists():
        return output_path

    input_path.write_text(mermaid_code, encoding="utf-8")

    cmd = [
        *mermaid_cli,
        "-i",
        str(input_path),
        "-o",
        str(output_path),
        "-b",
        "transparent",
        "-w",
        str(width),
        "-H",
        str(height),
    ]

    import os
    try:
        subprocess.run(cmd, check=True, capture_output=True, text=True, shell=(os.name == "nt"))
    except subprocess.CalledProcessError as exc:
        stderr = (exc.stderr or "").strip()
        print("  WARNING: Mermaid render failed; leaving Mermaid code block as text")
        if stderr:
            print(f"           Mermaid CLI: {stderr.splitlines()[-1]}")
        return None
    except OSError as exc:
        print(f"  WARNING: Mermaid render unavailable: {exc}")
        return None

    if output_path.exists():
        return output_path

    print("  WARNING: Mermaid render completed without output image")
    return None


def ensure_markdown_slide_entry(slide_path: str | Path) -> None:
    """Fail fast when a manifest slide entry is not a markdown source file."""
    path = Path(slide_path)
    ext = path.suffix.lower()
    if ext and ext not in ALLOWED_SLIDE_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_SLIDE_EXTENSIONS))
        raise ValueError(
            f"Non-markdown slide entry detected: {slide_path} "
            f"(extension: {ext}). Expected one of: {allowed}. "
            "Use a markdown deck source in sections[].decks, or configure a .pptx under the top-level template field."
        )


def extract_slide_title(file_path: Path) -> str:
    """Return the first ## H2 heading text, or the file stem as fallback."""
    ensure_markdown_slide_entry(file_path)
    try:
        text = file_path.read_text(encoding="utf-8-sig")
    except FileNotFoundError:
        return file_path.stem

    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            text = text[end + 4 :]

    for line in text.splitlines():
        m = re.match(r"^## (.+)", line)
        if m:
            return m.group(1).strip()
    return file_path.stem


def resolve_repo_path(repo_root: Path, candidate: str | Path) -> Path:
    """Resolve a manifest path relative to the repository root."""
    path = Path(candidate)
    if path.is_absolute():
        return path
    return repo_root / path


def get_slide_layout(prs: Presentation, preferred_names: list[str], fallback_index: int):
    """Return the first layout whose name matches, else fall back to an index."""
    normalized_names = {name.strip().casefold() for name in preferred_names if name.strip()}

    for layout in prs.slide_layouts:
        layout_name = getattr(layout, "name", "").strip().casefold()
        if layout_name in normalized_names:
            return layout

    try:
        return prs.slide_layouts[fallback_index]
    except IndexError:
        return prs.slide_layouts[0]


def find_slide_layout_by_name(prs: Presentation, layout_name: str):
    """Return the slide layout whose name matches the requested name, else None."""
    normalized_name = layout_name.strip().casefold()
    if not normalized_name:
        return None

    for layout in prs.slide_layouts:
        candidate_name = getattr(layout, "name", "").strip().casefold()
        if candidate_name == normalized_name:
            return layout

    return None


def format_available_layout_names(prs: Presentation) -> str:
    """Return a readable list of layout names available in the template."""
    names = []
    for index, layout in enumerate(prs.slide_layouts):
        layout_name = getattr(layout, "name", "").strip()
        names.append(layout_name or f"<unnamed:{index}>")
    return ", ".join(names)


def build_section_id(section_name: str, slide_start_idx: int) -> str:
    """Return a deterministic section id for the PPTX XML."""
    digest = hashlib.sha1(f"{section_name}:{slide_start_idx}".encode("utf-8")).hexdigest()
    return str(int(digest[:8], 16))


def split_marp_slides(md_content: str) -> list[str]:
    """Split a Marp markdown file into individual slide blocks.

    Strips the YAML front matter, then splits the remaining content on bare
    `---` lines that are NOT inside fenced code blocks.
    Returns a list of slide content strings (one per slide).
    """
    if md_content.startswith('\ufeff'):
        md_content = md_content[1:]
    lines = md_content.splitlines()

    # Strip YAML front matter. We can't just take the first closing `---`
    # because AI provenance `prompt: |` blocks may legitimately contain `---`
    # lines as plain scalar content.
    if lines and lines[0].strip() == "---":
        for end in range(1, len(lines)):
            if lines[end].strip() != "---":
                continue
            front_matter = "\n".join(lines[1:end])
            try:
                parsed = yaml.safe_load(front_matter)
            except yaml.YAMLError:
                continue
            if isinstance(parsed, dict):
                lines = lines[end + 1:]
                break
            else:
                print(f"Warning: Front matter parsed to {type(parsed)}, expected dict.")
                break
        else:
            print("Warning: Front matter stripping failed to find a valid closing ---.")
            pass

    # Split on bare --- separators outside fenced code blocks
    slides: list[list[str]] = []
    current: list[str] = []
    in_fence = False
    fence_pat = re.compile(r"^(`{3,}|~{3,})")

    for line in lines:
        if fence_pat.match(line):
            in_fence = not in_fence
        if line.strip() == "---" and not in_fence:
            slides.append(current)
            current = []
        else:
            current.append(line)
    slides.append(current)

    # Filter out empty or whitespace-only blocks
    return ["\n".join(block).strip() for block in slides if any(line.strip() for line in block)]


def process_markdown_links(text: str) -> str:
    """
    Process markdown links according to rendering rules:
    - If link text == link target: render just the link text
    - If link text != link target: render as "link text: link target"
    """
    def replace_link(match):
        link_text = match.group(1)
        link_target = match.group(2)

        if link_text == link_target:
            return link_text
        else:
            return f"{link_text}: {link_target}"

    # Match markdown links: [text](url)
    # Exclude markdown image tokens: ![alt](url)
    # Use non-greedy matching to handle multiple links on one line
    pattern = r'(?<!!)\[([^\]]+)\]\(([^\)]+)\)'
    return re.sub(pattern, replace_link, text)


def apply_markdown_formatting(text_frame, line_text: str, paragraph=None, force_monospace: bool = False) -> None:
    """
    Parse markdown inline formatting and blockquote syntax into PowerPoint runs.
    Supports **bold**, *italic*, _italic_, ~~strikethrough~~, <u>underline</u>, and
    inline code spans using backticks.

    Blockquote lines (starting with '> ') are rendered as italic text — PPTX has no
    native blockquote element so italic is the closest visual equivalent.

    Example: "> 'Programming hasn't changed'"
    -> Run 1: "'Programming hasn't changed'" (italic)
    """
    def set_bullet_visibility(paragraph, show_bullet: bool) -> None:
        p_pr = paragraph._p.get_or_add_pPr()
        for child in list(p_pr):
            if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
                p_pr.remove(child)
        if not show_bullet:
            p_pr.insert(0, OxmlElement("a:buNone"))

    if paragraph is None:
        p = text_frame.add_paragraph()
    else:
        p = paragraph
        p.clear()

    if force_monospace:
        p_pr = p._p.get_or_add_pPr()
        for child in list(p_pr):
            if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
                p_pr.remove(child)
        p_pr.insert(0, OxmlElement("a:buNone"))

        run = p.add_run()
        run.text = line_text
        run.font.name = "Consolas"
        return

    # Handle blockquote lines: render as italic, strip the '> ' prefix
    is_blockquote = line_text.startswith("> ")
    if is_blockquote:
        line_text = line_text[2:]

    # Preserve leading indentation so nested markdown lists map to PPT levels.
    leading_spaces = len(line_text) - len(line_text.lstrip(" "))
    list_level = min(leading_spaces // 2, 8)

    # Support task list markers: - [ ] item, - [x] item
    task_match = re.match(r"^\s*[-*+]\s+\[([ xX])\]\s+(.*)$", line_text)
    if task_match:
        checked = task_match.group(1).lower() == "x"
        line_text = f"[{'x' if checked else ' '}] {task_match.group(2)}"
        p.level = list_level
        set_bullet_visibility(p, False)
    # Support ordered list markers: 1. item
    elif re.match(r"^\s*\d+\.\s+", line_text):
        p.level = list_level
        set_bullet_visibility(p, False)
    # Support unordered list markers -, *, +
    bullet_match = re.match(r"^\s*([-*+])\s+(.*)$", line_text)
    if bullet_match and not task_match:
        line_text = bullet_match.group(2)
        p.level = list_level
        set_bullet_visibility(p, True)
    elif not task_match and not re.match(r"^\s*\d+\.\s+", line_text):
        set_bullet_visibility(p, False)

    token_pattern = re.compile(r"(<u>[^<]+</u>|\*\*[^*]+\*\*|~~[^~]+~~|`[^`\n]+`|\*[^*\n]+\*|_[^_\n]+_)")
    segments = []
    last_pos = 0
    for match in token_pattern.finditer(line_text):
        if match.start() > last_pos:
            segments.append((line_text[last_pos:match.start()], False, False, False, False, False))

        token = match.group(0)
        is_bold = token.startswith("**") and token.endswith("**")
        is_strike = token.startswith("~~") and token.endswith("~~")
        is_underline = token.startswith("<u>") and token.endswith("</u>")
        is_code = token.startswith("`") and token.endswith("`")
        is_italic = (
            (token.startswith("*") and token.endswith("*") and not is_bold)
            or (token.startswith("_") and token.endswith("_") and not (token.startswith("__") and token.endswith("__")))
        )

        if is_bold:
            inner = token[2:-2]
        elif is_strike:
            inner = token[2:-2]
        elif is_underline:
            inner = token[3:-4]
        elif is_code:
            inner = token[1:-1]
        elif is_italic:
            inner = token[1:-1]
        else:
            inner = token

        segments.append((inner, is_bold, is_italic, is_underline, is_strike, is_code))
        last_pos = match.end()

    if last_pos < len(line_text):
        segments.append((line_text[last_pos:], False, False, False, False, False))

    if not segments:
        segments = [(line_text, False, False, False, False, False)]

    for seg_text, seg_bold, seg_italic, seg_underline, seg_strike, seg_code in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.bold = seg_bold
        run.font.italic = seg_italic or is_blockquote
        run.font.underline = seg_underline
        run.font.strike = seg_strike
        if seg_strike:
            # Ensure reliable strike rendering across PowerPoint clients.
            run._r.get_or_add_rPr().set("strike", "sngStrike")
        if seg_code:
            run.font.name = "Consolas"


def contains_markdown_table(body: str) -> bool:
    """Check if the body contains a markdown table."""
    lines = body.strip().splitlines()
    for i, line in enumerate(lines):
        # A markdown table has a separator line with |---|---| pattern
        if re.match(r'^\s*\|?\s*[-:]+\s*\|', line):
            return True
    return False


def parse_markdown_table(body: str) -> list[list[str]]:
    """Parse a markdown table and return rows of cells.

    Returns a list of rows, where each row is a list of cell values.
    Example:
        | Header 1 | Header 2 |
        |----------|----------|
        | Cell 1   | Cell 2   |

    Returns: [['Header 1', 'Header 2'], ['Cell 1', 'Cell 2']]
    """
    lines = body.strip().splitlines()
    rows = []

    for line in lines:
        line = line.strip()
        # Skip separator lines (|---|---|)
        if re.match(r'^\|?\s*[-:]+\s*\|', line):
            continue
        # Parse table rows
        if '|' in line:
            # Remove leading/trailing pipes if present
            if line.startswith('|'):
                line = line[1:]
            if line.endswith('|'):
                line = line[:-1]
            # Split by pipe and strip whitespace
            cells = [cell.strip() for cell in line.split('|')]
            rows.append(cells)

    return rows


def parse_slide(md_content: str) -> tuple[str, str, str | None, str, str, list[str], bool, str]:
    """Return (title, body, background_image_path, speaker_notes, layout_name, inline_images, title_is_h1, matter_of_fact_title) parsed from a single markdown slide block.

    title_is_h1 is True when the slide title was set from a bare `# H1` heading
    rather than a `## H2` heading.  The caller uses this flag to decide whether
    to inject a Centered Title (Section Header) slide instead of a Title Only slide.

    When the H1 heading uses the ``matter of fact || witty`` convention:
    - ``title`` is set to the witty portion (after ``||``)
    - ``matter_of_fact_title`` is set to the descriptive portion (before ``||``)
    For all other slides ``matter_of_fact_title`` is an empty string.
    """
    lines = md_content.strip().splitlines()

    title = ""
    matter_of_fact_title = ""
    title_is_h1 = False
    body_lines = []
    bg_image = None
    speaker_notes = ""
    layout_name = ""
    inline_images: list[str] = []

    # Pattern to match Marp background images: ![bg ...](path)
    bg_pattern = re.compile(r'!\[bg[^\]]*\]\(([^)]+)\)')
    inline_image_pattern = re.compile(r'^!\[[^\]]*\]\(([^)]+)\)$')
    legacy_image_marker_pattern = re.compile(r'^!Slide\s+\d+\s+image:\s+(.+)$', re.IGNORECASE)
    layout_pattern = re.compile(r'^<!--\s*layout\s*:\s*(.*?)\s*-->$', re.IGNORECASE)

    for line in lines:
        if not title and line.startswith("# "):
            raw = line[2:].strip()
            # "Matter of Fact Title || Witty Title" — split into two parts
            if "||" in raw:
                parts = raw.split("||", 1)
                matter_of_fact_title = parts[0].strip()
                title = parts[1].strip()
            else:
                title = raw
            title_is_h1 = True
        elif line.startswith("## "):
            if not title:
                title = line[3:].strip()
        else:
            # Extract background image path if present
            stripped = line.strip()
            if stripped.startswith("![bg"):
                match = bg_pattern.match(stripped)
                if match:
                    bg_image = match.group(1)
            elif stripped.startswith("!["):
                match = inline_image_pattern.match(stripped)
                if match:
                    inline_images.append(match.group(1).strip())
            elif stripped.startswith("!Slide"):
                match = legacy_image_marker_pattern.match(stripped)
                if match:
                    inline_images.append(match.group(1).strip())
            elif stripped.startswith("<!--"):
                match = layout_pattern.match(stripped)
                if match:
                    layout_name = match.group(1).strip()
            # Skip Marp directives (background images and HTML comments)
            elif not stripped.startswith("<!--"):
                body_lines.append(line)

    # Extract speaker notes (::: notes ... :::) from body
    body = "\n".join(body_lines).strip()
    notes_pattern = re.compile(r":::[ \t]*notes(.*?):::", re.DOTALL | re.IGNORECASE)
    notes_match = notes_pattern.search(body)
    if notes_match:
        speaker_notes = notes_match.group(1).strip()
        body = notes_pattern.sub("", body).strip()

    # Process markdown links according to rendering rules
    body = process_markdown_links(body)

    return title, body, bg_image, speaker_notes, layout_name, inline_images, title_is_h1, matter_of_fact_title


def add_inline_images(slide, image_paths: list[str]) -> None:
    """Add inline images to the slide near the lower-right area.

    This keeps markdown image directives from being rendered as literal text.
    """
    if not image_paths:
        return

    base_top = Inches(3.6)
    gap = Inches(0.2)
    max_width = Inches(3.2)
    right_margin = Inches(0.4)

    current_top = base_top
    for image_path in image_paths:
        try:
            pic = slide.shapes.add_picture(str(image_path), left=0, top=0, width=max_width)
            pic.left = Inches(10) - right_margin - pic.width
            pic.top = current_top
            current_top = pic.top + pic.height + gap
        except Exception as e:
            print(f"  WARNING: Could not add inline image {image_path}: {e}")


def set_slide_notes(slide, note_text: str) -> None:
    """Set the speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


def populate_text_placeholder(shape, body: str) -> None:
    """Populate a text placeholder using the existing markdown formatting rules."""
    tf = shape.text_frame
    tf.clear()
    lines = body.splitlines()
    # Avoid an empty first paragraph when slide body starts with blank lines.
    while lines and not lines[0].strip():
        lines.pop(0)
    if lines:
        fence_pattern = re.compile(r"^(`{3,}|~{3,})")
        in_code_fence = False

        def add_line(line: str, paragraph=None) -> None:
            nonlocal in_code_fence
            stripped = line.strip()
            fence_match = fence_pattern.match(stripped)
            if fence_match:
                if not in_code_fence:
                    # Opening fence — suppress backticks and language label.
                    in_code_fence = True
                else:
                    # Closing fence — skip rendering entirely
                    in_code_fence = False
            else:
                apply_markdown_formatting(tf, line, paragraph=paragraph, force_monospace=in_code_fence)

        add_line(lines[0], paragraph=tf.paragraphs[0])
        for line in lines[1:]:
            add_line(line)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Save run-level font names (including None/inherited) before any prefit.
    # fit_text() can stamp a single font across all runs; we must restore each
    # run's original font intent so only explicitly marked code remains monospace.
    font_names_by_run: dict[tuple[int, int], str | None] = {}
    for pi, para in enumerate(tf.paragraphs):
        for ri, run in enumerate(para.runs):
            font_names_by_run[(pi, ri)] = run.font.name

    def restore_font_overrides() -> None:
        for pi, para in enumerate(tf.paragraphs):
            for ri, run in enumerate(para.runs):
                if (pi, ri) in font_names_by_run:
                    run.font.name = font_names_by_run[(pi, ri)]

    # Pre-fit once so slides open already fitted without requiring manual
    # toggle/reflow in PowerPoint. Use Consolas metrics when code runs exist
    # because it is wider than Arial and gives a safer initial fit.
    if font_names_by_run:
        has_consolas = any((name or "").lower() == "consolas" for name in font_names_by_run.values())
        prefit_family = "Consolas" if has_consolas else "Arial"
        prefit_max_size = 18 if has_consolas else 28
        try:
            tf.fit_text(font_family=prefit_family, max_size=prefit_max_size)
        except Exception:
            pass
        restore_font_overrides()

    finalize_text_frame(tf)

    # Restore run-level font overrides (e.g. Consolas for code blocks/spans).
    restore_font_overrides()


def populate_paragraph_lines(text_frame, lines: list[str]) -> None:
    """Populate a text frame without leaving an empty first paragraph behind."""
    text_frame.clear()
    if not lines:
        return

    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.text = lines[0]
    for line in lines[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line

    finalize_text_frame(text_frame)


def finalize_text_frame(text_frame, max_size: int = 28) -> None:
    """Configure text frame for live PowerPoint autofit behavior."""
    text_frame.word_wrap = True
    # Keep native autofit enabled so Office computes sizing with the actual
    # run fonts (e.g. Consolas in code blocks), not a precomputed Arial fit.
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Enforce OOXML autofit explicitly so template-level noAutofit/spAutoFit
    # settings cannot override the desired behavior in PowerPoint.
    body_pr = text_frame._txBody.bodyPr
    for child in list(body_pr):
        if child.tag.endswith("}noAutofit") or child.tag.endswith("}spAutoFit") or child.tag.endswith("}normAutofit"):
            body_pr.remove(child)
    body_pr.append(OxmlElement("a:normAutofit"))


def enforce_autofit_in_presentation(prs: Presentation) -> None:
    """Force autofit mode on every text frame in the presentation.

    Some templates/layout placeholders can reintroduce noAutofit through
    inheritance; this pass makes the final slide XML deterministic.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            body_pr = tf._txBody.bodyPr
            for child in list(body_pr):
                if child.tag.endswith("}noAutofit") or child.tag.endswith("}spAutoFit") or child.tag.endswith("}normAutofit"):
                    body_pr.remove(child)
            body_pr.append(OxmlElement("a:normAutofit"))


def style_table_medium_accent_4(table) -> None:
    """Apply a visual style matching PowerPoint's Medium Style 2 - Accent 4."""
    table.first_row = True
    table.horz_banding = True

    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.text_frame.word_wrap = True
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(16)

            if row_index == 0:
                cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
                paragraph.font.bold = True
                paragraph.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
            else:
                cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
                cell.fill.fore_color.brightness = 0.8 if row_index % 2 == 1 else 0.92
                paragraph.font.bold = False
                paragraph.font.color.theme_color = MSO_THEME_COLOR.TEXT_1


def populate_table_cell_markdown(cell, text: str, font_size: int = 16) -> None:
    """Populate a table cell using markdown inline formatting rules."""
    tf = cell.text_frame
    tf.clear()
    apply_markdown_formatting(tf, text, paragraph=tf.paragraphs[0])
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def split_two_column_body(body: str) -> tuple[str, str]:
    """Split slide body into left/right column text.

    Supported formats:
    - Explicit separator line: `::: column`
    - Two or more `###` subsections, where the first two become left/right columns
    """
    separator_pattern = re.compile(r"^\s*:::\s*column\s*$", re.IGNORECASE | re.MULTILINE)
    if separator_pattern.search(body):
        left, right = separator_pattern.split(body, maxsplit=1)
        return left.strip(), right.strip()

    sections: list[tuple[str | None, list[str]]] = []
    current_heading: str | None = None
    current_lines: list[str] = []

    for line in body.splitlines():
        if line.startswith("### "):
            if current_heading is not None or current_lines:
                sections.append((current_heading, current_lines))
            current_heading = line[4:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_heading is not None or current_lines:
        sections.append((current_heading, current_lines))

    if len(sections) >= 2:
        def render_section(section_heading: str | None, section_lines: list[str]) -> str:
            rendered_lines: list[str] = []
            if section_heading:
                rendered_lines.append(f"**{section_heading}**")
            rendered_lines.extend(section_lines)
            return "\n".join(rendered_lines).strip()

        left_heading, left_lines = sections[0]
        right_heading, right_lines = sections[1]
        return render_section(left_heading, left_lines), render_section(right_heading, right_lines)

    # Exercise fallback: split at Activities when Objectives/Activities labels are present.
    lines = body.splitlines()
    activity_idx = None
    has_objectives = False
    for idx, line in enumerate(lines):
        stripped = line.strip().lower()
        if stripped == "objectives:":
            has_objectives = True
        elif stripped == "activities:":
            activity_idx = idx
            break

    if has_objectives and activity_idx is not None and activity_idx > 0:
        left = "\n".join(lines[:activity_idx]).strip()
        right = "\n".join(lines[activity_idx:]).strip()
        if right:
            return left, right

    return body.strip(), ""


def has_explicit_two_column_separator(body: str) -> bool:
    """Return True when the slide body explicitly declares a two-column split."""
    separator_pattern = re.compile(r"^\s*:::\s*column\s*$", re.IGNORECASE | re.MULTILINE)
    return separator_pattern.search(body) is not None


def should_use_two_column_layout(title: str, body: str) -> bool:
    """Decide when to use two-column layout automatically.

    Uses explicit separators first, then applies an exercise-specific fallback.
    """
    if has_explicit_two_column_separator(body):
        return True

    if title.strip().lower().startswith("exercise:"):
        _, right = split_two_column_body(body)
        return bool(right)

    return False


def add_title_content_slide(prs: Presentation, title: str, body: str, inline_images: list[str] | None = None, note: str = "") -> None:
    """Add a Title and Content layout slide with markdown bold formatting support."""
    layout = get_slide_layout(prs, ["Title and Content"], LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            populate_text_placeholder(shape, body)
            break

    add_inline_images(slide, inline_images or [])

    if note:
        set_slide_notes(slide, note)


def add_background_image(slide, bg_image: str | None) -> None:
    """Add a full-slide background image behind existing content."""
    if not bg_image:
        return

    try:
        pic = slide.shapes.add_picture(
            bg_image,
            left=0,
            top=0,
            width=Inches(10),
            height=Inches(7.5)
        )
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    except Exception as e:
        print(f"  WARNING: Could not add background image {bg_image}: {e}")


def add_named_layout_slide(
    prs: Presentation,
    layout_name: str,
    title: str,
    body: str,
    inline_images: list[str] | None = None,
    bg_image: str | None = None,
    note: str = "",
) -> bool:
    """Add a slide using the exact layout name requested in slide markdown or manifest."""
    layout = find_slide_layout_by_name(prs, layout_name)
    if layout is None:
        print(
            f'  WARNING: layout "{layout_name}" not found in template. '
            f"Available layouts: {format_available_layout_names(prs)}"
        )
        return False

    slide = prs.slides.add_slide(layout)
    add_background_image(slide, bg_image)

    if slide.shapes.title:
        slide.shapes.title.text = title

    title_placeholder_idx = slide.shapes.title.placeholder_format.idx if slide.shapes.title else None
    content_placeholders = []
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == title_placeholder_idx:
            continue
        content_placeholders.append(shape)

    content_placeholders.sort(key=lambda shape: shape.placeholder_format.idx)

    left_body, right_body = split_two_column_body(body)
    if len(content_placeholders) >= 2 and right_body:
        populate_text_placeholder(content_placeholders[0], left_body)
        populate_text_placeholder(content_placeholders[1], right_body)
    elif content_placeholders and body:
        populate_text_placeholder(content_placeholders[0], body)

    add_inline_images(slide, inline_images or [])

    if note:
        set_slide_notes(slide, note)

    return True


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_body: str,
    right_body: str,
    inline_images: list[str] | None = None,
    note: str = "",
) -> None:
    """Add a two-column slide using the template's two-column layout."""
    layout = get_slide_layout(prs, ["Two Column", "Two Columns", "Two Content"], LAYOUT_TWO_COLUMN)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 1 and left_body:
            populate_text_placeholder(shape, left_body)
        elif idx == 2 and right_body:
            populate_text_placeholder(shape, right_body)

    add_inline_images(slide, inline_images or [])

    if note:
        set_slide_notes(slide, note)


def add_table_slide(prs: Presentation, title: str, table_data: list[list[str]], note: str = "") -> None:
    """Add a Title and Content layout slide with a table."""
    layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    # Remove the content placeholder to make room for the table
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            sp = shape.element
            sp.getparent().remove(sp)
            break

    # Calculate table dimensions
    rows = len(table_data)
    cols = max(len(row) for row in table_data) if table_data else 0

    if rows > 0 and cols > 0:
        # Position and size the table
        left = Inches(0.5)
        top = Inches(2.0)
        width = Inches(9.0)
        height = Inches(0.5 * rows)  # Dynamic height based on number of rows

        # Add table
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        style_table_medium_accent_4(table)

        # Populate table cells
        for i, row_data in enumerate(table_data):
            for j, cell_value in enumerate(row_data):
                if j < cols:  # Ensure we don't exceed column count
                    cell = table.cell(i, j)
                    populate_table_cell_markdown(cell, cell_value, font_size=16)

    if note:
        set_slide_notes(slide, note)


def add_title_only_slide(prs: Presentation, title: str, bg_image: str | None = None, inline_images: list[str] | None = None, note: str = "") -> None:
    """Add a Title Only layout slide (no body content), optionally with a background image."""
    try:
        layout = get_slide_layout(prs, ["Title Only"], LAYOUT_TITLE_ONLY)
    except IndexError:
        layout = get_slide_layout(prs, ["Title and Content"], LAYOUT_TITLE_CONTENT)

    slide = prs.slides.add_slide(layout)

    # Add background image first (before setting title) so it appears behind content
    add_background_image(slide, bg_image)

    # Set title after adding background image
    if slide.shapes.title:
        slide.shapes.title.text = title

    add_inline_images(slide, inline_images or [])

    if note:
        set_slide_notes(slide, note)


def add_title_slide(prs: Presentation, title: str, subtitle: str = "", note: str = "") -> None:
    """Add a Title Slide layout (typically used for presentation opening)."""
    try:
        layout = get_slide_layout(prs, ["Title Slide"], LAYOUT_TITLE_SLIDE)
    except IndexError:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # Title Slide typically has two placeholders: title (idx 0) and subtitle (idx 1)
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:
            shape.text = title
        elif idx == 1 and subtitle:
            shape.text = subtitle

    if note:
        set_slide_notes(slide, note)


def add_module_list_slide(prs: Presentation, all_sections: list[str], current: str, note: str = "") -> None:
    """Add a Course Modules navigation slide; current section is bold + arrow."""
    layout = get_slide_layout(prs, ["Title and Content"], LAYOUT_TITLE_CONTENT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Course Modules"

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            lines = [f"- **▶ {name}**" if name == current else f"- {name}" for name in all_sections]
            if lines:
                apply_markdown_formatting(tf, lines[0], paragraph=tf.paragraphs[0])
                for line in lines[1:]:
                    apply_markdown_formatting(tf, line)
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            break

    if note:
        set_slide_notes(slide, note)


def add_centered_title_slide(prs: Presentation, title: str, note: str = "") -> None:
    """Add a Section Header (centered title) slide for H1 deck-cover headings.

    This is injected for every `# H1` slide block after the very first one in the
    first deck of the manifest.  It uses the Section Header layout so the title
    is centred on a plain background, visually signalling a new deck topic.
    """
    layout = get_slide_layout(prs, ["Section Header", "Section Title"], LAYOUT_SECTION_HEADER)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if note:
        set_slide_notes(slide, note)


def add_centered_two_titles_slide(prs: Presentation, title: str, subtitle: str, note: str = "") -> None:
    """Add a Centered Two Titles slide for H1 deck-cover headings that use the ``matter of fact || witty`` convention.

    The descriptive (matter-of-fact) portion is placed in the Title placeholder.
    The witty portion is placed in the Subtitle placeholder below it.
    Both are centred on a plain background.
    """
    layout = get_slide_layout(prs, ["Centered Two Titles"], LAYOUT_CENTERED_TWO_TITLES)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    # Write the witty portion into the subtitle placeholder (idx=1)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = subtitle
            break
    if note:
        set_slide_notes(slide, note)


def build_presentation(yaml_path: Path, output_path: Path) -> None:
    yaml_path = yaml_path.resolve()
    repo_root = yaml_path.parent.parent

    with open(yaml_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    sections_cfg = config.get("sections", [])
    all_section_names = [s.get("name", "Unnamed Section") for s in sections_cfg]

    template = config.get("template")
    if template:
        template_path = resolve_repo_path(repo_root, template)
        if not template_path.exists():
            template_path = Path(template)
        prs = Presentation(str(template_path))
        print(f"Using template: {template_path}")
    else:
        prs = Presentation()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    mermaid_cli = detect_mermaid_cli()
    mermaid_cache_dir = output_path.parent / ".mermaid-cache"
    mermaid_warning_shown = False

    # Track first H1 across all decks so we can inject Centered Title slides.
    # The very first H1 encountered in the first deck is kept as-is (it is the
    # deck's own cover/title slide).  Every subsequent H1-only slide block is
    # replaced with a Section Header (Centered Title) slide.
    first_h1_seen = False
    global_deck_idx = 0  # Incremented once per deck file processed

    # Populate the template title slide (slide 0) if title/subtitle are provided.
    prs_title = config.get("title", "")
    prs_subtitle = config.get("subtitle", "")
    if prs.slides and (prs_title or prs_subtitle):
        title_slide = prs.slides[0]
        for shape in title_slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 0 and prs_title:
                shape.text = prs_title
            elif idx == 1 and prs_subtitle:
                shape.text = prs_subtitle

    # Populate the template Agenda slide (slide 1) with section names.
    if len(prs.slides) > 1:
        agenda_slide = prs.slides[1]
        # Only populate if title is "Agenda" (or placeholder is empty)
        agenda_title = ""
        for shape in agenda_slide.shapes:
            if shape.has_text_frame and "title" in shape.name.lower():
                agenda_title = shape.text_frame.text.strip()
                break
        if agenda_title.lower() == "agenda":
            for shape in agenda_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    populate_paragraph_lines(shape.text_frame, all_section_names)
                    break

    # python-pptx exposes the underlying <p:presentation> element via part._element.
    prs_el = prs.part._element
    ns14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    etree.register_namespace("p14", ns14)
    section_lst_tag = f"{{{ns14}}}sectionLst"
    section_lst = prs_el.find(section_lst_tag)
    if section_lst is not None:
        prs_el.remove(section_lst)
    section_lst = etree.SubElement(prs_el, section_lst_tag)

    for idx, section in enumerate(sections_cfg):
        section_name = section.get("name", "Unnamed Section")
        deck_entries = [s for s in (section.get("decks") or []) if s]

        slide_start_idx = len(prs.slides)

        if idx > 0:
            add_module_list_slide(
                prs,
                all_section_names,
                section_name,
                note="Auto-generated course navigation slide showing all modules with current section highlighted"
            )

        is_first_deck = (global_deck_idx == 0)

        for slide_entry in deck_entries:
            if isinstance(slide_entry, dict):
                if "layout" in slide_entry:
                    raise ValueError(
                        "Manifest deck layouts are no longer supported. "
                        "Move the layout into the slide markdown using "
                        '`<!-- layout: Layout Name -->`.'
                    )
                slide_path = slide_entry.get("file") or slide_entry.get("path")
            else:
                slide_path = slide_entry

            ensure_markdown_slide_entry(slide_path)

            slide_file = resolve_repo_path(repo_root, slide_path)
            if not slide_file.exists():
                print(f"  WARNING: not found -- {slide_file}")
                global_deck_idx += 1
                continue

            note = f"Source: {slide_path}"
            md = slide_file.read_text(encoding="utf-8-sig")
            slide_blocks = split_marp_slides(md)
            for block in slide_blocks:
                title, body, bg_image, speaker_notes, slide_layout_name, inline_image_refs, title_is_h1, matter_of_fact_title = parse_slide(block)
                rendered_mermaid_path: str | None = None

                body_without_mermaid, mermaid_blocks = extract_mermaid_blocks(body)
                if mermaid_blocks:
                    if mermaid_cli is None:
                        if not mermaid_warning_shown:
                            print(
                                "  WARNING: Mermaid CLI not found (install mmdc or npx). "
                                "Mermaid blocks will remain as code in PPTX text content."
                            )
                            mermaid_warning_shown = True
                    else:
                        if len(mermaid_blocks) > 1:
                            print("  WARNING: Multiple Mermaid blocks detected; rendering the first block only")

                        rendered_mermaid = render_mermaid_png(
                            mermaid_blocks[0],
                            mermaid_cache_dir,
                            mermaid_cli,
                        )
                        if rendered_mermaid is not None:
                            rendered_mermaid_path = str(rendered_mermaid)
                            body = body_without_mermaid

                # Combine source path with speaker notes if they exist
                if speaker_notes:
                    combined_note = f"{speaker_notes}\n\n---\n\nSource: {slide_path}"
                else:
                    combined_note = note

                # Inject a Centered Title slide for H1-only blocks after the
                # first H1 in the first deck.
                if title_is_h1 and not body and not slide_layout_name:
                    if not first_h1_seen and is_first_deck:
                        # This is the very first H1 in the first deck — keep as-is.
                        first_h1_seen = True
                    else:
                        first_h1_seen = True
                        if matter_of_fact_title:
                            # "Matter of Fact || Witty" — use Centered Two Titles layout.
                            # matter_of_fact_title → Title placeholder
                            # title (the witty part)  → Subtitle placeholder
                            add_centered_two_titles_slide(
                                prs,
                                title=matter_of_fact_title,
                                subtitle=title or slide_file.stem,
                                note=combined_note,
                            )
                        else:
                            add_centered_title_slide(
                                prs,
                                title or slide_file.stem,
                                note=combined_note,
                            )
                        continue

                # Resolve background image path relative to source deck first,
                # then fall back to manifest folder and repo root.
                bg_image_path = None
                if bg_image:
                    bg_image_path = (slide_file.parent / bg_image).resolve()
                    if not bg_image_path.exists():
                        bg_image_path = (yaml_path.parent / bg_image).resolve()
                    if not bg_image_path.exists():
                        bg_image_path = resolve_repo_path(repo_root, bg_image)
                    if not bg_image_path.exists():
                        print(f"  WARNING: background image not found -- {bg_image_path}")
                        bg_image_path = None
                    else:
                        bg_image_path = str(bg_image_path)

                inline_image_paths: list[str] = []
                for image_ref in inline_image_refs:
                    image_path = (slide_file.parent / image_ref).resolve()
                    if not image_path.exists():
                        image_path = (yaml_path.parent / image_ref).resolve()
                    if not image_path.exists():
                        image_path = resolve_repo_path(repo_root, image_ref)
                    if not image_path.exists():
                        print(f"  WARNING: inline image not found -- {image_path}")
                        continue
                    inline_image_paths.append(str(image_path))

                if rendered_mermaid_path is not None:
                    inline_image_paths.append(rendered_mermaid_path)

                requested_layout_name = slide_layout_name

                if requested_layout_name and add_named_layout_slide(
                    prs,
                    requested_layout_name,
                    title or slide_file.stem,
                    body,
                    inline_images=inline_image_paths,
                    bg_image=bg_image_path,
                    note=combined_note,
                ):
                    continue

                if body and should_use_two_column_layout(title or slide_file.stem, body):
                    left_body, right_body = split_two_column_body(body)
                    if right_body:
                        add_two_column_slide(
                            prs,
                            title or slide_file.stem,
                            left_body,
                            right_body,
                            inline_images=inline_image_paths,
                            note=combined_note,
                        )
                        continue

                if contains_markdown_table(body):
                    # Parse and create table slide
                    table_data = parse_markdown_table(body)
                    add_table_slide(prs, title or slide_file.stem, table_data, note=combined_note)
                elif body:
                    add_title_content_slide(
                        prs,
                        title or slide_file.stem,
                        body,
                        inline_images=inline_image_paths,
                        note=combined_note,
                    )
                else:
                    add_title_only_slide(
                        prs,
                        title or slide_file.stem,
                        bg_image=bg_image_path,
                        inline_images=inline_image_paths,
                        note=combined_note,
                    )

        slide_end_idx = len(prs.slides)
        all_sld_ids = list(prs.slides._sldIdLst)
        section_el = etree.SubElement(
            section_lst,
            f"{{{ns14}}}section",
            attrib={
                "name": section_name,
                "id": build_section_id(section_name, slide_start_idx),
            },
        )
        sld_id_lst_section = etree.SubElement(section_el, f"{{{ns14}}}sldIdLst")
        for sld_el in all_sld_ids[slide_start_idx:slide_end_idx]:
            etree.SubElement(
                sld_id_lst_section,
                f"{{{ns14}}}sldId",
                attrib={"id": sld_el.get("id")},
            )

            global_deck_idx += 1

        if not deck_entries:
            print(f"  INFO: Section '{section_name}' is empty — only module list slide added")

    enforce_autofit_in_presentation(prs)
    prs.save(output_path)
    print(f"✅ Saved: {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build PPTX from YAML manifest")
    parser.add_argument("yaml_path", help="Path to the YAML manifest file")
    parser.add_argument("output_path", help="Path for the generated PPTX")
    args = parser.parse_args()
    build_presentation(Path(args.yaml_path), Path(args.output_path))
