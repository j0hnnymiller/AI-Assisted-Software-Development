"""
generate_pptx.py -- Build a PPTX from a YAML manifest using python-pptx.
Requires: pip install python-pptx pyyaml

Usage:
    python generate_pptx.py <yaml_path> <output_pptx_path>
"""
import argparse
import hashlib
import re
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation

LAYOUT_TITLE_CONTENT = 1  # adjust index if template differs
LAYOUT_SECTION_HEADER = 2  # adjust index if template differs
LAYOUT_TITLE_ONLY = 5  # adjust index if template differs


def extract_slide_title(file_path: Path) -> str:
    """Return the first ## H2 heading text, or the file stem as fallback."""
    try:
        text = file_path.read_text(encoding="utf-8")
    except FileNotFoundError:
        return file_path.stem

    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            text = text[end + 4 :]

    for line in text.splitlines():
        m = re.match(r"^## (.+)", line)
        if m:
            return m.group(1).strip()
    return file_path.stem


def resolve_repo_path(repo_root: Path, candidate: str | Path) -> Path:
    """Resolve a manifest path relative to the repository root."""
    path = Path(candidate)
    if path.is_absolute():
        return path
    return repo_root / path


def build_section_id(section_name: str, slide_start_idx: int) -> str:
    """Return a deterministic section id for the PPTX XML."""
    digest = hashlib.sha1(f"{section_name}:{slide_start_idx}".encode("utf-8")).hexdigest()
    return str(int(digest[:8], 16))


def split_marp_slides(md_content: str) -> list[str]:
    """Split a Marp markdown file into individual slide blocks.

    Strips the YAML front matter (first --- ... --- block), then splits the
    remaining content on bare `---` lines that are NOT inside fenced code blocks.
    Returns a list of slide content strings (one per slide).
    """
    lines = md_content.splitlines()

    # Strip YAML front matter
    if lines and lines[0].strip() == "---":
        end = next(
            (i for i, line in enumerate(lines[1:], 1) if line.strip() == "---"),
            None,
        )
        if end is not None:
            lines = lines[end + 1:]

    # Split on bare --- separators outside fenced code blocks
    slides: list[list[str]] = []
    current: list[str] = []
    in_fence = False
    fence_pat = re.compile(r"^(`{3,}|~{3,})")

    for line in lines:
        if fence_pat.match(line):
            in_fence = not in_fence
        if line.strip() == "---" and not in_fence:
            slides.append(current)
            current = []
        else:
            current.append(line)
    slides.append(current)

    # Filter out empty or whitespace-only blocks
    return ["\n".join(block).strip() for block in slides if any(l.strip() for l in block)]


def parse_slide(md_content: str) -> tuple[str, str]:
    """Return (title, body) parsed from a single markdown slide block."""
    lines = md_content.strip().splitlines()

    title = ""
    body_lines = []
    for line in lines:
        if not title and line.startswith("# "):
            title = line[2:].strip()
        elif line.startswith("## "):
            if not title:
                title = line[3:].strip()
        else:
            body_lines.append(line)

    # Strip speaker notes (::: notes ... :::) from body
    body = "\n".join(body_lines).strip()
    body = re.sub(r":::[ \t]*notes.*?:::", "", body, flags=re.DOTALL).strip()
    return title, body


def set_slide_notes(slide, note_text: str) -> None:
    """Set the speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


def add_title_content_slide(prs: Presentation, title: str, body: str, note: str = "") -> None:
    """Add a Title and Content layout slide."""
    layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            for line in body.splitlines():
                p = tf.add_paragraph()
                if line.startswith("- "):
                    p.text = line[2:]
                    p.level = 0
                else:
                    p.text = line
            break

    if note:
        set_slide_notes(slide, note)


def add_title_only_slide(prs: Presentation, title: str, note: str = "") -> None:
    """Add a Title Only layout slide (no body content)."""
    try:
        layout = prs.slide_layouts[LAYOUT_TITLE_ONLY]
    except IndexError:
        layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]

    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    if note:
        set_slide_notes(slide, note)


def add_section_header_slide(prs: Presentation, section_name: str) -> None:
    """Add a section header slide (# Section Name, lead layout)."""
    try:
        layout = prs.slide_layouts[LAYOUT_SECTION_HEADER]
    except IndexError:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = section_name


def add_module_list_slide(prs: Presentation, all_sections: list[str], current: str) -> None:
    """Add a Course Modules navigation slide; current section is bold + arrow."""
    layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Course Modules"

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            for name in all_sections:
                p = tf.add_paragraph()
                p.level = 0
                run = p.add_run()
                if name == current:
                    run.text = f"▶ {name}"
                    run.font.bold = True
                else:
                    run.text = name
            break


def add_section_agenda_slide(
    prs: Presentation, section_name: str, slide_files: list, repo_root: Path
) -> None:
    """Add an agenda slide listing titles extracted from source files."""
    if not slide_files:
        return

    titles = [extract_slide_title(resolve_repo_path(repo_root, f)) for f in slide_files]
    add_title_content_slide(prs, section_name, "\n".join(f"- {t}" for t in titles))


def build_presentation(yaml_path: Path, output_path: Path) -> None:
    yaml_path = yaml_path.resolve()
    repo_root = yaml_path.parent.parent

    with open(yaml_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    sections_cfg = config.get("sections", [])
    all_section_names = [s.get("name", "Unnamed Section") for s in sections_cfg]

    template = config.get("template")
    if template:
        template_path = resolve_repo_path(repo_root, template)
        if not template_path.exists():
            template_path = Path(template)
        prs = Presentation(str(template_path))
        print(f"Using template: {template_path}")
    else:
        prs = Presentation()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Populate the template title slide (slide 0) if title/subtitle are provided.
    prs_title = config.get("title", "")
    prs_subtitle = config.get("subtitle", "")
    if prs.slides and (prs_title or prs_subtitle):
        title_slide = prs.slides[0]
        for shape in title_slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 0 and prs_title:
                shape.text = prs_title
            elif idx == 1 and prs_subtitle:
                shape.text = prs_subtitle

    # Populate the template Agenda slide (slide 1) with section names.
    if len(prs.slides) > 1:
        agenda_slide = prs.slides[1]
        # Only populate if title is "Agenda" (or placeholder is empty)
        agenda_title = ""
        for shape in agenda_slide.shapes:
            if shape.has_text_frame and "title" in shape.name.lower():
                agenda_title = shape.text_frame.text.strip()
                break
        if agenda_title.lower() == "agenda":
            for shape in agenda_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    tf = shape.text_frame
                    tf.clear()
                    for name in all_section_names:
                        p = tf.add_paragraph()
                        p.text = name
                    break

    # python-pptx exposes the underlying <p:presentation> element via part._element.
    prs_el = prs.part._element
    ns14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    etree.register_namespace("p14", ns14)
    section_lst_tag = f"{{{ns14}}}sectionLst"
    section_lst = prs_el.find(section_lst_tag)
    if section_lst is not None:
        prs_el.remove(section_lst)
    section_lst = etree.SubElement(prs_el, section_lst_tag)

    for section in sections_cfg:
        section_name = section.get("name", "Unnamed Section")
        slides = [s for s in (section.get("slides") or []) if s]

        slide_start_idx = len(prs.slides)

        add_module_list_slide(prs, all_section_names, section_name)
        add_section_header_slide(prs, section_name)
        add_section_agenda_slide(prs, section_name, slides, repo_root)

        for slide_path in slides:
            slide_file = resolve_repo_path(repo_root, slide_path)
            if not slide_file.exists():
                print(f"  WARNING: not found -- {slide_file}")
                continue

            note = f"Source: {slide_path}"
            md = slide_file.read_text(encoding="utf-8")
            slide_blocks = split_marp_slides(md)
            for block in slide_blocks:
                title, body = parse_slide(block)
                if body:
                    add_title_content_slide(prs, title or slide_file.stem, body, note=note)
                else:
                    add_title_only_slide(prs, title or slide_file.stem, note=note)

        slide_end_idx = len(prs.slides)
        all_sld_ids = list(prs.slides._sldIdLst)
        section_el = etree.SubElement(
            section_lst,
            f"{{{ns14}}}section",
            attrib={
                "name": section_name,
                "id": build_section_id(section_name, slide_start_idx),
            },
        )
        sld_id_lst_section = etree.SubElement(section_el, f"{{{ns14}}}sldIdLst")
        for sld_el in all_sld_ids[slide_start_idx:slide_end_idx]:
            etree.SubElement(
                sld_id_lst_section,
                f"{{{ns14}}}sldId",
                attrib={"id": sld_el.get("id")},
            )

        if not slides:
            print(f"  INFO: Section '{section_name}' is empty — only injected slides added")

    prs.save(output_path)
    print(f"✅ Saved: {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build PPTX from YAML manifest")
    parser.add_argument("yaml_path", help="Path to the YAML manifest file")
    parser.add_argument("output_path", help="Path for the generated PPTX")
    args = parser.parse_args()
    build_presentation(Path(args.yaml_path), Path(args.output_path))
