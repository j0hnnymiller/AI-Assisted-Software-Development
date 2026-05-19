#!/usr/bin/env python3
"""
Extract PPTX to Marp Markdown format with speaker notes.
Requires: pip install python-pptx
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    sys.exit(1)

TITLE_PLACEHOLDER_TYPES = {1, 13, 15}   # TITLE, CENTER_TITLE, and legacy values
SUBTITLE_PLACEHOLDER_TYPES = {2, 12}    # SUBTITLE


def get_placeholder_type(shape):
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type
    except Exception:
        pass
    return None


def normalize_markdown_line(level, text):
    """Normalize a paragraph into a markdown-safe line preserving indentation."""
    indent = "  " * level
    if text.startswith(("•", "○", "▪", "-", "*", "–")):
        return f"{indent}- {text[1:].strip()}"
    if level > 0:
        return f"{indent}- {text}"
    return text


def extract_text_from_shape(shape):
    """Extract text from a non-title shape, including tables."""
    text_parts = []
    try:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    level = paragraph.level
                    text_parts.append(normalize_markdown_line(level, text))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                if rows:
                    text_parts.append(" | ".join(rows[0]))
                    text_parts.append(" | ".join(["---"] * len(rows[0])))
                    for row in rows[1:]:
                        text_parts.append(" | ".join(row))
            except Exception:
                pass
    except Exception:
        pass
    return text_parts


def save_shape_image(shape, images_dir, pptx_stem, slide_num):
    """Persist a picture shape to disk and return saved filename, or None on failure."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None

    try:
        image = shape.image
        ext = image.ext
        filename = f"{pptx_stem}_slide{slide_num:02d}_{shape.shape_id}.{ext}"
        (images_dir / filename).write_bytes(image.blob)
        print(f"  Saved image: {filename}")
        return filename
    except Exception as exc:
        print(f"  Warning: could not save image on slide {slide_num}: {exc}")
        return None


def convert_pptx_to_marp(pptx_path, output_path):
    """Convert PPTX to Marp Markdown with ::: notes speaker notes."""
    pptx_path = Path(pptx_path)
    output_path = Path(output_path)

    if not pptx_path.exists():
        print(f"Error: Input file not found: {pptx_path}")
        return False

    print(f"Loading presentation: {pptx_path}")
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        print(f"Error opening PPTX: {exc}")
        return False

    output_path.parent.mkdir(parents=True, exist_ok=True)
    images_dir = output_path.parent / "images"
    images_dir.mkdir(exist_ok=True)
    pptx_stem = pptx_path.stem.replace(" ", "_")

    lines = ["---", "marp: true", "theme: default", "paginate: true", "---", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {slide_num}...")

        title_text = None
        body_lines = []
        image_files = []

        for shape in slide.shapes:
            ph_type = get_placeholder_type(shape)
            if ph_type in TITLE_PLACEHOLDER_TYPES or ph_type in SUBTITLE_PLACEHOLDER_TYPES:
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        if title_text is None:
                            title_text = t
                        else:
                            body_lines.append(normalize_markdown_line(0, t))
            else:
                body_lines.extend(extract_text_from_shape(shape))

            img_filename = save_shape_image(shape, images_dir, pptx_stem, slide_num)
            if img_filename:
                image_files.append(img_filename)

        if title_text:
            lines.append(f"## {title_text}")
        if body_lines:
            lines.append("")
            lines.extend(body_lines)

        for img_filename in image_files:
            lines.append("")
            lines.append(f"![Slide {slide_num} image](images/{img_filename})")

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("")
                lines.append("::: notes")
                for note_line in notes_text.splitlines():
                    lines.append(note_line.strip())
                lines.append(":::")

        lines.append("")
        lines.append("---")
        lines.append("")

        print(
            f"  Slide {slide_num}: {len(body_lines)} text lines, "
            f"{len(image_files)} image(s), "
            f"{'notes' if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip() else 'no notes'}"
        )

    print(f"Writing output to: {output_path}")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"Successfully converted {len(prs.slides)} slides")
    return True


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_to_marp.py <input.pptx> [output.md]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) >= 3 else str(Path(input_path).with_suffix(".md"))

    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)

    ok = convert_pptx_to_marp(input_path, output_path)
    sys.exit(0 if ok else 1)
